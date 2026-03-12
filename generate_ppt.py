from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height

# ── Color Palette ──
NAVY      = RGBColor(10, 25, 49)
DARK_BLUE = RGBColor(0, 63, 135)
BLUE      = RGBColor(0, 102, 204)
LIGHT_BLUE= RGBColor(224, 237, 255)
YELLOW    = RGBColor(255, 204, 0)
GOLD      = RGBColor(255, 183, 0)
WHITE     = RGBColor(255, 255, 255)
OFF_WHITE = RGBColor(245, 247, 252)
BLACK     = RGBColor(25, 25, 35)
GRAY      = RGBColor(120, 130, 145)
LIGHT_GRAY= RGBColor(210, 215, 225)
GREEN     = RGBColor(34, 166, 90)
RED       = RGBColor(220, 60, 50)
ORANGE    = RGBColor(245, 155, 30)

def _rect(slide, l, t, w, h, fill_rgb, line=False):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    if not line:
        s.line.fill.background()
    return s

def _rounded(slide, l, t, w, h, fill_rgb, line_rgb=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        s.line.color.rgb = line_rgb; s.line.width = Pt(1.5)
    else:
        s.line.fill.background()
    return s

def _text(slide, l, t, w, h, txt, sz=14, bold=False, color=BLACK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True; tf.auto_size = None
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(sz); p.font.bold = bold
    p.font.color.rgb = color; p.alignment = align
    tf.vertical_anchor = anchor
    return box

def _multiline(slide, l, t, w, h, lines, sz=14, color=BLACK, spacing=6, bold_first=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz); p.font.color.rgb = color
        p.space_after = Pt(spacing)
        if bold_first and i == 0:
            p.font.bold = True
    return box

def _line(slide, x1, y1, x2, y2, color_rgb=GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    ln = connector.line
    ln.color.rgb = color_rgb
    ln.width = width

def _arrow(slide, x1, y1, x2, y2, color_rgb=GRAY, width=Pt(2)):
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)  # straight
    ln = connector.line
    ln.color.rgb = color_rgb
    ln.width = width
    # Add arrowhead
    ln_xml = ln._ln
    tail = ln_xml.makeelement('{http://schemas.openxmlformats.org/drawingml/2006/main}tailEnd', {})
    tail.set('type', 'triangle')
    tail.set('w', 'med')
    tail.set('len', 'med')
    ln_xml.append(tail)

def slide_bg(slide, color=OFF_WHITE):
    bg = slide.background; f = bg.fill; f.solid(); f.fore_color.rgb = color

def title_band(slide, title, subtitle=None):
    """Top header band with title"""
    _rect(slide, Inches(0), Inches(0), SW, Inches(1.05), NAVY)
    _rect(slide, Inches(0), Inches(1.05), SW, Inches(0.05), YELLOW)
    _text(slide, Inches(0.7), Inches(0.15), Inches(10), Inches(0.7),
          title, sz=28, bold=True, color=WHITE, anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        _text(slide, Inches(0.7), Inches(0.6), Inches(10), Inches(0.4),
              subtitle, sz=13, color=YELLOW, anchor=MSO_ANCHOR.TOP)
    # Slide number area
    _text(slide, Inches(11.5), Inches(0.25), Inches(1.5), Inches(0.5),
          "TechTrack", sz=11, color=RGBColor(100,130,180), align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)

def info_card(slide, l, t, w, h, title, body_lines, accent=BLUE):
    """Rounded card with accent top strip"""
    card = _rounded(slide, l, t, w, h, WHITE, line_rgb=LIGHT_GRAY)
    _rect(slide, l + Inches(0.05), t + Inches(0.03), w - Inches(0.1), Inches(0.045), accent)
    _text(slide, l + Inches(0.2), t + Inches(0.15), w - Inches(0.4), Inches(0.35),
          title, sz=14, bold=True, color=accent)
    _multiline(slide, l + Inches(0.2), t + Inches(0.5), w - Inches(0.4), h - Inches(0.6),
               body_lines, sz=11, color=BLACK, spacing=4)

def erd_box(slide, l, t, w, title, fields, color=BLUE):
    """ERD entity box with title bar and field list"""
    row_h = Inches(0.22)
    h = Inches(0.38) + row_h * len(fields)
    # Title header
    header = _rect(slide, l, t, w, Inches(0.38), color)
    tf = header.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = title; p.font.size = Pt(10); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Fields
    body = _rect(slide, l, t + Inches(0.38), w, row_h * len(fields), WHITE)
    body.line.color.rgb = color; body.line.width = Pt(1.5)
    box = slide.shapes.add_textbox(l + Inches(0.08), t + Inches(0.4), w - Inches(0.16), row_h * len(fields))
    tf2 = box.text_frame; tf2.word_wrap = True
    for i, fld in enumerate(fields):
        p = tf2.paragraphs[0] if i == 0 else tf2.add_paragraph()
        p.text = fld; p.font.size = Pt(8); p.font.color.rgb = BLACK
        p.space_after = Pt(1)
        if fld.startswith("PK") or fld.startswith("FK"):
            p.font.bold = True
    return t + Inches(0.38) + row_h * len(fields)

def flow_box(slide, l, t, w, h, txt, fill=BLUE, txtcolor=WHITE, shape_type=MSO_SHAPE.ROUNDED_RECTANGLE):
    s = slide.shapes.add_shape(shape_type, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(9); p.font.bold = True
    p.font.color.rgb = txtcolor; p.alignment = PP_ALIGN.CENTER
    return s

def diamond(slide, l, t, w, h, txt, fill=GOLD):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    s.fill.solid(); s.fill.fore_color.rgb = fill
    s.line.fill.background()
    tf = s.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = txt; p.font.size = Pt(8); p.font.bold = True
    p.font.color.rgb = BLACK; p.alignment = PP_ALIGN.CENTER
    return s

# ════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s1, NAVY)

# Decorative shapes
_rect(s1, Inches(0), Inches(0), Inches(0.15), SH, YELLOW)
_rect(s1, Inches(0.15), Inches(0), Inches(0.04), SH, GOLD)

_text(s1, Inches(1.5), Inches(1.5), Inches(10), Inches(1.2),
      "TechTrack", sz=60, bold=True, color=WHITE, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.BOTTOM)

_rect(s1, Inches(1.5), Inches(2.85), Inches(3.5), Inches(0.05), YELLOW)

_text(s1, Inches(1.5), Inches(3.1), Inches(10), Inches(0.8),
      "Hardware Monitoring & Maintenance System", sz=22, color=YELLOW, align=PP_ALIGN.LEFT)

_multiline(s1, Inches(1.5), Inches(4.3), Inches(8), Inches(1.2), [
    "Java Swing Desktop Application",
    "Built with OOP Fundamentals  •  No External Dependencies",
    "Final Version"
], sz=14, color=RGBColor(160, 180, 220), spacing=6)

# Bottom bar
_rect(s1, Inches(0), Inches(6.9), SW, Inches(0.6), DARK_BLUE)
_text(s1, Inches(1.5), Inches(6.95), Inches(10), Inches(0.45),
      "CS / IT Project Presentation  •  March 2026", sz=12, color=RGBColor(140,160,200))

# ════════════════════════════════════════════════════════════════════
# SLIDE 2 — SYSTEM OVERVIEW
# ════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s2)
title_band(s2, "System Overview", "What is TechTrack?")

_multiline(s2, Inches(0.7), Inches(1.4), Inches(6), Inches(1.5), [
    "TechTrack is a desktop application for monitoring hardware",
    "equipment, scheduling maintenance, logging repairs, and",
    "managing issue reports — all in one system."
], sz=15, color=BLACK, spacing=8)

# Feature highlight cards (2 rows of 3)
cards = [
    ("Login & Auth", ["Secure username/password login", "Session management", "Change password"], BLUE),
    ("Dashboard", ["Live summary statistics", "Maintenance intelligence", "Activity log"], DARK_BLUE),
    ("Equipment", ["CRUD operations", "Search & filter", "Health score tracking"], GREEN),
    ("Issue Reports", ["Report issues by equipment", "Priority levels", "Status workflow"], ORANGE),
    ("Repair Logs", ["Log repairs with costs", "Track technicians", "Mark completions"], RED),
    ("Maintenance", ["Schedule tasks", "Due date tracking", "Status management"], RGBColor(100, 60, 180)),
]

for i, (title, lines, accent) in enumerate(cards):
    col = i % 3
    row = i // 3
    x = Inches(0.7) + col * Inches(4.1)
    y = Inches(3.2) + row * Inches(2.0)
    info_card(s2, x, y, Inches(3.8), Inches(1.8), title, lines, accent)

# ════════════════════════════════════════════════════════════════════
# SLIDE 3 — CODE STRUCTURE
# ════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s3)
title_band(s3, "Code Structure", "3-Layer Architecture")

# Architecture diagram — 3 vertical columns
layers = [
    ("UI LAYER (Presentation)", BLUE, [
        "LoginFrame.java", "MainDashboard.java",
        "DashboardPage.java", "EquipmentPage.java",
        "IssueReportsPage.java", "RepairLogsPage.java",
        "MaintenancePage.java", "SessionManager.java",
        "LogoUtils.java"
    ]),
    ("SERVICE LAYER (Business Logic)", GREEN, [
        "Manageable<T>  (Interface)",
        "DataStore  (Singleton)",
        "AuthService",
        "EquipmentService",
        "IssueReportService",
        "RepairLogService",
        "MaintenanceService"
    ]),
    ("MODEL LAYER (Data)", ORANGE, [
        "BaseEntity  (Abstract Class)",
        "Equipment",
        "IssueReport",
        "RepairLog",
        "MaintenanceTask",
        "User"
    ]),
]

for i, (lbl, accent, files) in enumerate(layers):
    x = Inches(0.5) + i * Inches(4.2)
    y = Inches(1.5)
    # Layer header
    hdr = _rounded(s3, x, y, Inches(3.9), Inches(0.5), accent)
    tf = hdr.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = lbl; p.font.size = Pt(13); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # Files as rounded cards
    fy = y + Inches(0.65)
    for f in files:
        card = _rounded(s3, x + Inches(0.1), fy, Inches(3.7), Inches(0.35), WHITE, line_rgb=LIGHT_GRAY)
        tf2 = card.text_frame; tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
        p2 = tf2.paragraphs[0]; p2.text = "   " + f; p2.font.size = Pt(10); p2.font.color.rgb = BLACK
        p2.alignment = PP_ALIGN.LEFT
        fy += Inches(0.38)

# Arrows between layers
for i in range(2):
    ax = Inches(0.5) + (i+1) * Inches(4.2) - Inches(0.15)
    _arrow(s3, ax - Inches(0.3), Inches(3.2), ax + Inches(0.1), Inches(3.2), GRAY, Pt(2.5))
    _arrow(s3, ax + Inches(0.1), Inches(3.5), ax - Inches(0.3), Inches(3.5), GRAY, Pt(2.5))

_text(s3, Inches(3.6), Inches(3.0), Inches(1.5), Inches(0.3), "calls ▸", sz=8, color=GRAY, align=PP_ALIGN.CENTER)
_text(s3, Inches(3.6), Inches(3.55), Inches(1.5), Inches(0.3), "◂ returns", sz=8, color=GRAY, align=PP_ALIGN.CENTER)
_text(s3, Inches(7.8), Inches(3.0), Inches(1.5), Inches(0.3), "uses ▸", sz=8, color=GRAY, align=PP_ALIGN.CENTER)
_text(s3, Inches(7.8), Inches(3.55), Inches(1.5), Inches(0.3), "◂ data", sz=8, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 4 — OOP PILLARS
# ════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s4)
title_band(s4, "Four Pillars of OOP")

pillars = [
    ("Encapsulation", BLUE,
     "All model fields are private.\nAccess only through public\ngetters and setters.\n\nSessionManager hides\nsession state behind\nstatic methods."),
    ("Inheritance", GREEN,
     "Equipment, IssueReport,\nRepairLog, MaintenanceTask,\nand User all extend the\nBaseEntity abstract class.\n\nCommon fields (id, createdDate)\nare inherited automatically."),
    ("Abstraction", ORANGE,
     "BaseEntity defines abstract\nmethods: toTableRow() and\ngetDisplayName().\n\nManageable<T> interface\nhides CRUD implementation\ndetails from the UI layer."),
    ("Polymorphism", RGBColor(130, 50, 180),
     "Each service implements\nManageable<T> differently.\n\nEach model overrides\ntoTableRow() with its own\ncolumn data.\n\nSame interface, different\nbehavior per type."),
]

for i, (name, accent, desc) in enumerate(pillars):
    x = Inches(0.5) + i * Inches(3.15)
    y = Inches(1.6)
    # Card
    _rounded(s4, x, y, Inches(2.95), Inches(5.2), WHITE, line_rgb=LIGHT_GRAY)
    # Accent strip
    _rect(s4, x + Inches(0.05), y + Inches(0.03), Inches(2.85), Inches(0.05), accent)
    # Number circle
    circ = s4.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.1), y + Inches(0.25), Inches(0.7), Inches(0.7))
    circ.fill.solid(); circ.fill.fore_color.rgb = accent; circ.line.fill.background()
    tf = circ.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.text = str(i+1); p.font.size = Pt(22); p.font.bold = True
    p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
    # Title
    _text(s4, x + Inches(0.15), y + Inches(1.1), Inches(2.65), Inches(0.4),
          name, sz=16, bold=True, color=accent, align=PP_ALIGN.CENTER)
    # Description
    for j, line in enumerate(desc.split('\n')):
        _text(s4, x + Inches(0.2), y + Inches(1.6) + j * Inches(0.28), Inches(2.55), Inches(0.28),
              line, sz=10, color=BLACK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 5 — ERD (Entity Relationship Diagram)
# ════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s5)
title_band(s5, "Entity Relationship Diagram (ERD)")

# BaseEntity (parent)
base_x = Inches(5.2)
base_y = Inches(1.35)
erd_box(s5, base_x, base_y, Inches(2.8), "BaseEntity  «abstract»", [
    "PK  id : int",
    "    createdDate : String",
    "    toTableRow() : String[]  «abstract»",
    "    getDisplayName() : String  «abstract»"
], NAVY)

# Child entities
children = [
    (Inches(0.3), Inches(3.6), Inches(2.4), "User", [
        "PK  id : int",
        "    username : String",
        "    password : String",
        "    role : String"
    ], BLUE),
    (Inches(2.9), Inches(3.6), Inches(2.4), "Equipment", [
        "PK  id : int",
        "    deviceName : String",
        "    type : String",
        "    status : String",
        "    healthScore : int",
        "    location : String",
        "    category : String"
    ], GREEN),
    (Inches(5.5), Inches(3.6), Inches(2.4), "IssueReport", [
        "PK  id : int",
        "FK  equipmentName : String",
        "    description : String",
        "    priority : String",
        "    status : String",
        "    reportedBy : String"
    ], ORANGE),
    (Inches(8.1), Inches(3.6), Inches(2.4), "RepairLog", [
        "PK  id : int",
        "FK  equipmentName : String",
        "    description : String",
        "    technician : String",
        "    status : String",
        "    cost : double",
        "    dateCompleted : String"
    ], RED),
    (Inches(10.7), Inches(3.6), Inches(2.4), "MaintenanceTask", [
        "PK  id : int",
        "FK  equipmentName : String",
        "    task : String",
        "    dueDate : String",
        "    status : String"
    ], RGBColor(130, 50, 180)),
]

for cx, cy, cw, ctitle, cfields, ccolor in children:
    erd_box(s5, cx, cy, cw, ctitle, cfields, ccolor)

# Inheritance arrows (from each child up to BaseEntity)
for cx, cy, cw, ctitle, cfields, ccolor in children:
    child_center_x = cx + cw / 2
    parent_bottom_y = base_y + Inches(1.3)
    _arrow(s5, child_center_x, cy, child_center_x, parent_bottom_y, ccolor, Pt(1.5))

# "extends" label
_text(s5, Inches(5.5), Inches(2.85), Inches(2.5), Inches(0.3),
      "▲  extends (inheritance)", sz=9, bold=True, color=GRAY, align=PP_ALIGN.CENTER)

# Relationships — Equipment is referenced by IssueReport, RepairLog, Maintenance
_text(s5, Inches(3.4), Inches(6.8), Inches(8), Inches(0.4),
      "IssueReport, RepairLog, and MaintenanceTask reference Equipment by equipmentName  (1 Equipment → many Issues/Repairs/Tasks)",
      sz=10, color=GRAY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════
# SLIDE 6 — SYSTEM FLOWCHART
# ════════════════════════════════════════════════════════════════════
s6 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s6)
title_band(s6, "System Flowchart")

# --- Main vertical flow (Login) ---
cx = Inches(2.8)  # center x for left column
bw = Inches(1.8)
bh = Inches(0.5)

# Start
flow_box(s6, cx - bw/2, Inches(1.4), bw, bh, "START", NAVY, WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
_arrow(s6, cx, Inches(1.9), cx, Inches(2.15), NAVY)

# Login Screen
flow_box(s6, cx - bw/2, Inches(2.15), bw, bh, "Login Screen", BLUE, WHITE)
_arrow(s6, cx, Inches(2.65), cx, Inches(2.9), NAVY)

# Enter Credentials
flow_box(s6, cx - Inches(1.1), Inches(2.9), Inches(2.2), bh, "Enter Username\n& Password", LIGHT_BLUE, BLACK)
_arrow(s6, cx, Inches(3.4), cx, Inches(3.65), NAVY)

# Diamond — valid?
diamond(s6, cx - Inches(0.7), Inches(3.65), Inches(1.4), Inches(0.9), "Valid\nCredentials?")

# No — loop back
_text(s6, cx + Inches(0.75), Inches(3.85), Inches(0.5), Inches(0.3), "No", sz=9, bold=True, color=RED)
_arrow(s6, cx + Inches(0.7), Inches(4.1), cx + Inches(1.5), Inches(4.1), RED)
# Error box
flow_box(s6, cx + Inches(1.5), Inches(3.85), Inches(1.6), Inches(0.45), "Show Error", RED, WHITE)
_line(s6, cx + Inches(2.3), Inches(3.85), cx + Inches(2.3), Inches(2.4), RED)
_arrow(s6, cx + Inches(2.3), Inches(2.4), cx + bw/2, Inches(2.4), RED)

# Yes — go to dashboard
_text(s6, cx - Inches(0.15), Inches(4.55), Inches(0.5), Inches(0.3), "Yes", sz=9, bold=True, color=GREEN)
_arrow(s6, cx, Inches(4.55), cx, Inches(4.85), GREEN)

# Set Session
flow_box(s6, cx - bw/2, Inches(4.85), bw, bh, "Set Session", GREEN, WHITE)
_arrow(s6, cx, Inches(5.35), cx, Inches(5.6), NAVY)

# Main Dashboard
flow_box(s6, cx - Inches(1.1), Inches(5.6), Inches(2.2), Inches(0.55), "Main Dashboard", NAVY, WHITE)

# --- Right side: Navigation from Dashboard ---
dcx = Inches(8.5)  # right column center

# Navigation label
_text(s6, Inches(5.1), Inches(5.7), Inches(1.8), Inches(0.4), "Navigate ▸", sz=10, bold=True, color=GRAY, align=PP_ALIGN.CENTER)
_arrow(s6, cx + Inches(1.1), Inches(5.87), Inches(6.9), Inches(5.87), GRAY)

# The 5 module pages
modules = [
    ("Dashboard Page", DARK_BLUE),
    ("Equipment Page", GREEN),
    ("Issue Reports Page", ORANGE),
    ("Repair Logs Page", RED),
    ("Maintenance Page", RGBColor(130, 50, 180)),
]
module_bw = Inches(2.0)
module_bh = Inches(0.42)

for i, (mod_name, mod_color) in enumerate(modules):
    my = Inches(1.5) + i * Inches(0.58)
    flow_box(s6, dcx - module_bw/2, my, module_bw, module_bh, mod_name, mod_color, WHITE)

# Vertical distribution line from navigate endpoint up to module pages
nav_x = Inches(6.9)
first_mod_center = Inches(1.5) + module_bh / 2
last_mod_center = Inches(1.5) + 4 * Inches(0.58) + module_bh / 2
_line(s6, nav_x, first_mod_center, nav_x, Inches(5.87), GRAY, Pt(1.2))

# Arrows from vertical line to each module
for i in range(5):
    my = Inches(1.5) + i * Inches(0.58) + module_bh / 2
    _arrow(s6, nav_x, my, dcx - module_bw/2, my, GRAY, Pt(1.2))

# CRUD box on far right
crud_x = Inches(11.0)
crud_y = Inches(1.5)
_rounded(s6, crud_x, crud_y, Inches(2.0), Inches(3.5), WHITE, line_rgb=LIGHT_GRAY)
_rect(s6, crud_x + Inches(0.03), crud_y + Inches(0.03), Inches(1.94), Inches(0.04), BLUE)
_text(s6, crud_x + Inches(0.1), crud_y + Inches(0.15), Inches(1.8), Inches(0.3),
      "CRUD Operations", sz=11, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

crud_items = ["Add / Create", "View / Read", "Edit / Update", "Delete", "Search / Filter"]
for j, ci in enumerate(crud_items):
    cy = crud_y + Inches(0.55) + j * Inches(0.5)
    flow_box(s6, crud_x + Inches(0.2), cy, Inches(1.6), Inches(0.35), ci, LIGHT_BLUE, BLACK)

# Arrows from module pages to CRUD
_arrow(s6, dcx + module_bw/2, Inches(2.7), crud_x, Inches(2.7), GRAY, Pt(1.2))
_text(s6, Inches(9.7), Inches(2.4), Inches(1.2), Inches(0.3), "performs ▸", sz=8, color=GRAY, align=PP_ALIGN.CENTER)

# Logout flow
_text(s6, cx - Inches(2.2), Inches(5.65), Inches(1.0), Inches(0.4), "Logout", sz=10, bold=True, color=RED)
_line(s6, cx - Inches(1.1), Inches(5.87), cx - Inches(1.6), Inches(5.87), RED)
_line(s6, cx - Inches(1.6), Inches(5.87), cx - Inches(1.6), Inches(2.4), RED)
_arrow(s6, cx - Inches(1.6), Inches(2.4), cx - bw/2, Inches(2.4), RED)

# ════════════════════════════════════════════════════════════════════
# SLIDE 7 — FEATURES: LOGIN & DASHBOARD
# ════════════════════════════════════════════════════════════════════
s7 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s7)
title_band(s7, "Features: Login & Dashboard")

# Login card
info_card(s7, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), "Login System", [
    "• Username and password authentication",
    "• Credentials validated through AuthService",
    "• Show / hide password toggle checkbox",
    "• Session created via SessionManager on success",
    "• Invalid login shows error dialog",
    "• Logout clears session and returns to login",
    "",
    "Default credentials:",
    "  Username: admin",
    "  Password: password"
], BLUE)

# Dashboard card
info_card(s7, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.2), "Dashboard", [
    "Live Summary Cards (calculated from services):",
    "  • Total Devices — count of all equipment",
    "  • Online & Healthy — count by status",
    "  • Needs Attention — flagged devices",
    "  • Open Repair Tickets — incomplete repairs",
    "  • Avg Health Score — average across all",
    "",
    "Maintenance Intelligence:",
    "  • Ranks devices by repair frequency",
    "  • Shows category and failure rate",
    "",
    "Activity Log — session events and timestamps"
], DARK_BLUE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 8 — FEATURES: EQUIPMENT & ISSUES
# ════════════════════════════════════════════════════════════════════
s8 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s8)
title_band(s8, "Features: Equipment & Issue Reports")

info_card(s8, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), "Equipment Management", [
    "Full CRUD with form dialogs:",
    "  • Add — device name, type, status, health, location, category",
    "  • Edit — select row, modify via dialog",
    "  • Delete — confirmation before removal",
    "",
    "Search:",
    "  • Filter by device name, type, or status",
    "  • Real-time table refresh",
    "",
    "Types: Workstation, Server, Network Switch, Router, Printer",
    "Statuses: Online, Offline, Needs Attention",
    "Categories: Hardware, Software, Network"
], GREEN)

info_card(s8, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), "Issue Reports", [
    "Report issues linked to registered equipment:",
    "  • Select equipment from dropdown",
    "  • Enter description and set priority",
    "  • Auto-set reporter from current session",
    "",
    "Priority Levels: Low, Medium, High, Critical",
    "",
    "Status Workflow:",
    "  Open → In Progress → Resolved → Closed",
    "",
    "Actions:",
    "  • Report Issue — open submission form",
    "  • Resolve — mark as resolved",
    "  • Close Issue — mark as closed",
    "  • Delete — remove with confirmation",
    "  • Filter — dropdown filter by status"
], ORANGE)

# ════════════════════════════════════════════════════════════════════
# SLIDE 9 — FEATURES: REPAIRS & MAINTENANCE
# ════════════════════════════════════════════════════════════════════
s9 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s9)
title_band(s9, "Features: Repair Logs & Maintenance")

info_card(s9, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.3), "Repair Logs", [
    "Log repairs for equipment:",
    "  • Select equipment from dropdown",
    "  • Enter description, technician name, cost",
    "  • New repairs default to Pending status",
    "",
    "Status Workflow:",
    "  Pending → In Progress → Completed",
    "",
    "Mark Complete:",
    "  • Auto-fills completion date (today's date)",
    "  • Updates status to Completed",
    "",
    "Actions: Add Repair, Mark Complete, Delete",
    "Filter: dropdown by status (All / Pending / In Progress / Completed)"
], RED)

info_card(s9, Inches(6.8), Inches(1.5), Inches(6.0), Inches(5.3), "Maintenance Schedule", [
    "Schedule maintenance tasks:",
    "  • Select equipment from dropdown",
    "  • Enter task description and due date",
    "  • New tasks default to Scheduled status",
    "",
    "Status Workflow:",
    "  Scheduled → In Progress → Completed / Overdue",
    "",
    "Actions:",
    "  • Schedule Task — create new task",
    "  • Start — change status to In Progress",
    "  • Complete — mark as Completed",
    "  • Delete — remove with confirmation",
    "",
    "Filter: dropdown by status"
], RGBColor(130, 50, 180))

# ════════════════════════════════════════════════════════════════════
# SLIDE 10 — TECH STACK
# ════════════════════════════════════════════════════════════════════
s10 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s10)
title_band(s10, "Technology Stack")

tech_items = [
    ("Language", "Java (JDK 25)", BLUE),
    ("GUI", "Java Swing", GREEN),
    ("Build Tool", "Apache Maven", ORANGE),
    ("Data Storage", "In-memory (ArrayList)", RGBColor(130, 50, 180)),
    ("Design Pattern", "Singleton, MVC-inspired", DARK_BLUE),
    ("Dependencies", "None — pure Java only", RED),
]

for i, (label, value, accent) in enumerate(tech_items):
    row = i // 2
    col = i % 2
    x = Inches(0.7) + col * Inches(6.2)
    y = Inches(1.8) + row * Inches(1.7)
    card = _rounded(s10, x, y, Inches(5.8), Inches(1.4), WHITE, line_rgb=LIGHT_GRAY)
    _rect(s10, x + Inches(0.05), y + Inches(0.03), Inches(5.7), Inches(0.05), accent)
    _text(s10, x + Inches(0.3), y + Inches(0.2), Inches(5.2), Inches(0.35),
          label, sz=13, bold=True, color=accent)
    _text(s10, x + Inches(0.3), y + Inches(0.6), Inches(5.2), Inches(0.6),
          value, sz=18, bold=True, color=BLACK)

# ════════════════════════════════════════════════════════════════════
# SLIDE 11 — THANK YOU
# ════════════════════════════════════════════════════════════════════
s11 = prs.slides.add_slide(prs.slide_layouts[6])
slide_bg(s11, NAVY)

_rect(s11, Inches(0), Inches(0), Inches(0.15), SH, YELLOW)
_rect(s11, Inches(0.15), Inches(0), Inches(0.04), SH, GOLD)

_text(s11, Inches(1), Inches(2.0), Inches(11), Inches(1.5),
      "Thank You", sz=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

_rect(s11, Inches(5), Inches(3.7), Inches(3), Inches(0.05), YELLOW)

_text(s11, Inches(1), Inches(4.0), Inches(11), Inches(0.8),
      "TechTrack — Hardware Monitoring & Maintenance System", sz=18, color=YELLOW, align=PP_ALIGN.CENTER)

_text(s11, Inches(1), Inches(5.0), Inches(11), Inches(0.5),
      "Java  •  Swing  •  OOP  •  Maven", sz=14, color=RGBColor(140,160,200), align=PP_ALIGN.CENTER)

_rect(s11, Inches(0), Inches(6.9), SW, Inches(0.6), DARK_BLUE)

# ── SAVE ──
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "TechTrack_Presentation.pptx")
prs.save(output_path)
print(f"Saved: {output_path}")
